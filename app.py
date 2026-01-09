import os, io, base64, smtplib
from datetime import datetime
from flask import Flask, request, jsonify, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageOps, ImageDraw, ImageFont
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

app = Flask(__name__, static_folder='.', static_url_path='')

# =======配置区域=======
SMTP_SERVER = 'smtp.exmail.qq.com'
SMTP_PORT = 465
SENDER_EMAIL = 'tylerqin@gtiggs.com' 
SENDER_PASSWORD = 'vfm3x9Rd74PMCF56' 
# ====================

def process_img_data(b64):
    """解码并修正图片旋转"""
    if ',' in b64: b64 = b64.split(',')[1]
    img = Image.open(io.BytesIO(base64.b64decode(b64)))
    return ImageOps.exif_transpose(img)

def create_preview_image(text, images_b64):
    """生成一张模拟PPT布局的图片供预览"""
    # 1. 创建一张 16:9 的白底图片 (缩小尺寸以便网络传输, 比如 800x450)
    W, H = 800, 450
    canvas = Image.new('RGB', (W, H), 'white')
    draw = ImageDraw.Draw(canvas)
    
    # 2. 绘制顶部文字 (尝试加载字体，如果失败用默认)
    try:
        font = ImageFont.truetype("arial.ttf", 40)
    except:
        font = ImageFont.load_default()
        
    # 文字居中计算 (简化版)
    draw.text((W/2, 30), text, fill='black', anchor="mm", font=font)
    
    # 3. 绘制三张图片 (模拟 PPT 的下部排列)
    # PPT逻辑: 顶边距约 1/3 处, 三图并排
    y_start = 150
    img_w = 200 # 模拟宽度
    img_h = int(img_w * 4 / 3) # 3:4 比例
    gap = 30
    x_start = (W - (img_w * 3 + gap * 2)) // 2
    
    for i, b64 in enumerate(images_b64):
        img = process_img_data(b64)
        img.thumbnail((img_w, img_h)) # 缩放
        # 粘贴到画布
        canvas.paste(img, (x_start + i*(img_w+gap), y_start))
        
    output = io.BytesIO()
    canvas.save(output, format='JPEG', quality=80)
    return base64.b64encode(output.getvalue()).decode()

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/preview', methods=['POST'])
def preview():
    """只生成预览图，不发邮件"""
    try:
        data = request.json
        img_b64 = create_preview_image(data['text'], data['images'])
        return jsonify({"status": "success", "preview": "data:image/jpeg;base64," + img_b64})
    except Exception as e:
        return jsonify({"status": "error", "msg": str(e)}), 500

@app.route('/send', methods=['POST'])
def send_mail():
    """生成真正的PPT并发送"""
    try:
        data = request.json
        text = data['text']
        imgs = data['images']
        recipients = data['recipients'].split(',')
        
        # --- 1. 生成 PPT (同之前逻辑) ---
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        
        for i, img_b64 in enumerate(imgs):
            img = process_img_data(img_b64)
            output = io.BytesIO()
            img.save(output, format='JPEG')
            output.seek(0)
            slide.shapes.add_picture(output, Inches(1.66 + i*4), Inches(2.5), width=Inches(3.5))

        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        # --- 2. 发送邮件 ---
        filename = f"{text}_{datetime.now().strftime('%Y%m%d')}.pptx"
        msg = MIMEMultipart()
        msg['Subject'] = f"试身照片: {text}"
        msg['From'] = SENDER_EMAIL
        msg['To'] = ",".join(recipients)
        
        body = f"Hi All,\n附件是[{text}]的试身照片, 请注意查收. 谢谢.\n"
        msg.attach(MIMEText(body, 'plain'))
        
        part = MIMEBase('application', "octet-stream")
        part.set_payload(ppt_io.read())
        encoders.encode_base64(part)
        part.add_header('Content-Disposition', f'attachment; filename="{filename}"')
        msg.attach(part)
        
        with smtplib.SMTP_SSL(SMTP_SERVER, SMTP_PORT) as s:
            s.login(SENDER_EMAIL, SENDER_PASSWORD)
            s.sendmail(SENDER_EMAIL, recipients, msg.as_string())
            
        return jsonify({"status": "success"})
    except Exception as e:
        return jsonify({"status": "error", "msg": str(e)}), 500

if __name__ == '__main__':
    app.run()
